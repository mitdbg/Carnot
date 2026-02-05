from __future__ import annotations

import io
from typing import Any

import fsspec
import pandas as pd


class DataItem:
    """
    Reference to a data item in Carnot. Each item is stored at the provided absolute path.
    Supports local and S3 paths for various business file formats.
    """
    def __init__(self, path: str | None, embedding: list[float] | None = None, metadata: dict | None = None):
        self.path = path
        self.embedding = embedding
        self.metadata = metadata or {}
        self._dict = None

    def _read_file_contents(self) -> str:
        """
        Reads the file from self.path (local or S3) and parses its content based on extension.
        """
        if not self.path:
            return ""

        # use fsspec to open the file regardless of storage (s3:// or local)
        with fsspec.open(self.path, "rb") as f:
            suffix = self.path.split(".")[-1].lower()
            
            try:
                if suffix == "txt":
                    return f.read().decode("utf-8")
                
                elif suffix == "csv":
                    df = pd.read_csv(f)
                    return df.to_string()
                
                elif suffix == "parquet":
                    df = pd.read_parquet(f)
                    return df.to_string()
                
                elif suffix in ["xlsx", "xls"]:
                    df = pd.read_excel(f)
                    return df.to_string()
                
                elif suffix == "pdf":
                    import fitz
                    doc = fitz.open(stream=f.read(), filetype="pdf")
                    return "\n".join([page.get_text() for page in doc])
                
                elif suffix == "docx":
                    from docx import Document
                    doc = Document(io.BytesIO(f.read()))
                    return "\n".join([para.text for para in doc.paragraphs])
                
                elif suffix == "pptx":
                    from pptx import Presentation
                    pres = Presentation(io.BytesIO(f.read()))
                    text_runs = []
                    for slide in pres.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_runs.append(shape.text)
                    return "\n".join(text_runs)

                else:
                    # fallback for unknown file types
                    return f.read().decode("utf-8", errors="ignore")
            except Exception as e:
                return f"Error reading file {self.path}: {str(e)}"

    def update_dict(self, data: dict[str, Any]) -> None:
        """
        Updates only the internal cache dictionary. 
        Does not sync back to instance attributes (path, embedding, etc.).
        """
        if self._dict is None:
            self.to_dict()
        
        self._dict.update(data)

    def to_dict(self) -> dict:
        """
        Returns the internal dictionary. If it doesn't exist, it is initialized 
        from the current instance attributes and file contents.
        """
        if self._dict is None:
            contents = self._read_file_contents()
            self._dict = {
                "contents": contents, 
                "path": self.path, 
                "metadata": self.metadata,
                "embedding": self.embedding
            }
        return self._dict

    @staticmethod
    def from_dict(item_dict: dict, metadata_override: dict | None = None) -> DataItem:
        """
        Initializes a DataItem where path and embedding are None.
        The state is primarily held in the internal _dict attribute.
        """
        # path and embedding are initialized as None as requested
        instance = DataItem(path=None, embedding=None)
        
        # Only modify the metadata attribute if metadata_override is provided
        if metadata_override is not None:
            # Sync the metadata in the dict with the override
            current_metadata = item_dict.get("metadata", {})
            current_metadata.update(metadata_override)
            item_dict["metadata"] = current_metadata
            instance.metadata = current_metadata
            
        instance._dict = item_dict
        return instance
